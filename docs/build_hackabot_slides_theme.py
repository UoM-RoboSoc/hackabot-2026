from __future__ import annotations

import sys
import zipfile
from pathlib import Path
from tempfile import NamedTemporaryFile

sys.path.insert(0, "/tmp/codex-py")

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


DOCS_DIR = Path(__file__).resolve().parent
OUTPUT_PATH = DOCS_DIR / "hackabot-2026-google-slides-theme.pptx"
REFERENCE_IMAGE = DOCS_DIR / "theme-color-scheme-visual-light-board.png"

PALETTE = {
    "bg_1": "101113",
    "bg_2": "161922",
    "surface": "0F1318",
    "panel": "212730",
    "line": "2B2D42",
    "text": "F1F3F5",
    "text_dim": "C1C2C5",
    "text_soft": "D9DDE4",
    "accent": "EF233C",
    "accent_strong": "D90429",
    "accent_soft": "FFACB7",
    "cool_gray": "8D99AE",
    "neutral": "EDF2F4",
}

NS = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
}


def rgb(name_or_hex: str) -> RGBColor:
    value = PALETTE.get(name_or_hex, name_or_hex).replace("#", "")
    return RGBColor.from_string(value)


def set_shape_transparency(shape, value: float) -> None:
    shape.fill.transparency = value


def remove_line(shape) -> None:
    shape.line.fill.background()


def send_to_back(shape) -> None:
    sp = shape.element
    parent = sp.getparent()
    parent.remove(sp)
    parent.insert(2, sp)


def set_text(
    shape,
    text: str,
    *,
    font_name: str = "Space Grotesk",
    size: int = 18,
    color: str = "text",
    bold: bool = False,
    align=PP_ALIGN.LEFT,
    level: int = 0,
) -> None:
    text_frame = shape.text_frame
    text_frame.clear()
    text_frame.word_wrap = True
    text_frame.vertical_anchor = MSO_ANCHOR.TOP
    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = align
    paragraph.level = level
    run = paragraph.add_run()
    run.text = text
    font = run.font
    font.name = font_name
    font.size = Pt(size)
    font.bold = bold
    font.color.rgb = rgb(color)


def style_paragraph(
    paragraph,
    *,
    font_name: str,
    size: int,
    color: str,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
) -> None:
    paragraph.alignment = align
    for run in paragraph.runs:
        font = run.font
        font.name = font_name
        font.size = Pt(size)
        font.bold = bold
        font.color.rgb = rgb(color)


def style_placeholder(
    shape,
    *,
    font_name: str,
    size: int,
    color: str,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
) -> None:
    text_frame = shape.text_frame
    text_frame.word_wrap = True
    text_frame.vertical_anchor = MSO_ANCHOR.TOP
    for paragraph in text_frame.paragraphs:
        style_paragraph(
            paragraph,
            font_name=font_name,
            size=size,
            color=color,
            bold=bold,
            align=align,
        )


def add_glow(slide_like, left, top, width, height, color: str, transparency: float) -> None:
    shape = slide_like.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.OVAL,
        left,
        top,
        width,
        height,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(color)
    set_shape_transparency(shape, transparency)
    remove_line(shape)
    send_to_back(shape)


def add_top_rule(slide_like, width) -> None:
    shape = slide_like.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        0,
        0,
        width,
        Inches(0.08),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb("accent_strong")
    remove_line(shape)


def add_panel(slide_like, left, top, width, height, *, fill: str, transparency: float, line: str) -> None:
    shape = slide_like.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        left,
        top,
        width,
        height,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    set_shape_transparency(shape, transparency)
    shape.line.color.rgb = rgb(line)
    shape.line.transparency = 0.3
    send_to_back(shape)


def add_chip(
    slide_like,
    left,
    top,
    width,
    height,
    text: str,
    *,
    fill: str,
    text_color: str,
    line: str | None = None,
    transparency: float = 0.0,
) -> None:
    shape = slide_like.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        left,
        top,
        width,
        height,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    set_shape_transparency(shape, transparency)
    if line:
        shape.line.color.rgb = rgb(line)
    else:
        remove_line(shape)
    set_text(shape, text, size=11, color=text_color, bold=True, font_name="Space Grotesk", align=PP_ALIGN.CENTER)
    shape.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE


def add_label(
    slide_like,
    left,
    top,
    width,
    height,
    text: str,
    *,
    size: int,
    color: str,
    bold: bool = False,
    font_name: str = "Space Grotesk",
    align=PP_ALIGN.LEFT,
):
    text_box = slide_like.shapes.add_textbox(left, top, width, height)
    set_text(
        text_box,
        text,
        font_name=font_name,
        size=size,
        color=color,
        bold=bold,
        align=align,
    )
    return text_box


def style_layouts(prs: Presentation) -> None:
    master = prs.slide_masters[0]
    master.background.fill.solid()
    master.background.fill.fore_color.rgb = rgb("bg_1")

    for layout in master.slide_layouts:
        layout.background.fill.solid()
        layout.background.fill.fore_color.rgb = rgb("bg_1")

    title_layout = master.slide_layouts[0]
    title_layout.placeholders[0].left = Inches(0.9)
    title_layout.placeholders[0].top = Inches(1.3)
    title_layout.placeholders[0].width = Inches(8.3)
    title_layout.placeholders[0].height = Inches(1.4)
    style_placeholder(
        title_layout.placeholders[0],
        font_name="Space Grotesk",
        size=28,
        color="text",
        bold=True,
    )
    title_layout.placeholders[1].left = Inches(0.92)
    title_layout.placeholders[1].top = Inches(2.75)
    title_layout.placeholders[1].width = Inches(7.2)
    title_layout.placeholders[1].height = Inches(1.2)
    style_placeholder(
        title_layout.placeholders[1],
        font_name="Arial",
        size=14,
        color="text_dim",
    )

    content_layout = master.slide_layouts[1]
    content_layout.placeholders[0].left = Inches(0.75)
    content_layout.placeholders[0].top = Inches(0.55)
    content_layout.placeholders[0].width = Inches(10.7)
    content_layout.placeholders[0].height = Inches(0.75)
    content_layout.placeholders[1].left = Inches(0.95)
    content_layout.placeholders[1].top = Inches(1.75)
    content_layout.placeholders[1].width = Inches(10.8)
    content_layout.placeholders[1].height = Inches(4.7)
    style_placeholder(
        content_layout.placeholders[0],
        font_name="Space Grotesk",
        size=22,
        color="text",
        bold=True,
    )
    style_placeholder(
        content_layout.placeholders[1],
        font_name="Arial",
        size=18,
        color="text_dim",
    )

    section_layout = master.slide_layouts[2]
    section_layout.background.fill.solid()
    section_layout.background.fill.fore_color.rgb = rgb("bg_2")
    section_layout.placeholders[0].left = Inches(0.85)
    section_layout.placeholders[0].top = Inches(3.85)
    section_layout.placeholders[0].width = Inches(8.8)
    section_layout.placeholders[0].height = Inches(1.0)
    section_layout.placeholders[1].left = Inches(0.88)
    section_layout.placeholders[1].top = Inches(5.0)
    section_layout.placeholders[1].width = Inches(8.4)
    section_layout.placeholders[1].height = Inches(0.7)
    style_placeholder(
        section_layout.placeholders[0],
        font_name="Space Grotesk",
        size=26,
        color="text",
        bold=True,
    )
    style_placeholder(
        section_layout.placeholders[1],
        font_name="Arial",
        size=14,
        color="text_dim",
    )

    two_content_layout = master.slide_layouts[3]
    two_content_layout.placeholders[0].left = Inches(0.75)
    two_content_layout.placeholders[0].top = Inches(0.55)
    two_content_layout.placeholders[0].width = Inches(10.7)
    two_content_layout.placeholders[0].height = Inches(0.75)
    two_content_layout.placeholders[1].left = Inches(0.95)
    two_content_layout.placeholders[1].top = Inches(1.9)
    two_content_layout.placeholders[1].width = Inches(5.0)
    two_content_layout.placeholders[1].height = Inches(4.3)
    two_content_layout.placeholders[2].left = Inches(7.05)
    two_content_layout.placeholders[2].top = Inches(1.9)
    two_content_layout.placeholders[2].width = Inches(5.0)
    two_content_layout.placeholders[2].height = Inches(4.3)
    style_placeholder(
        two_content_layout.placeholders[0],
        font_name="Space Grotesk",
        size=22,
        color="text",
        bold=True,
    )
    style_placeholder(
        two_content_layout.placeholders[1],
        font_name="Arial",
        size=17,
        color="text_dim",
    )
    style_placeholder(
        two_content_layout.placeholders[2],
        font_name="Arial",
        size=17,
        color="text_dim",
    )

    title_only_layout = master.slide_layouts[5]
    title_only_layout.placeholders[0].left = Inches(0.75)
    title_only_layout.placeholders[0].top = Inches(0.55)
    title_only_layout.placeholders[0].width = Inches(11.0)
    title_only_layout.placeholders[0].height = Inches(0.75)
    style_placeholder(
        title_only_layout.placeholders[0],
        font_name="Space Grotesk",
        size=22,
        color="text",
        bold=True,
    )

    for layout in master.slide_layouts:
        for placeholder in layout.placeholders:
            placeholder_type = str(placeholder.placeholder_format.type)
            if placeholder_type in {"DATE (16)", "FOOTER (15)", "SLIDE_NUMBER (13)"}:
                if hasattr(placeholder, "text_frame"):
                    style_placeholder(
                        placeholder,
                        font_name="Arial",
                        size=9,
                        color="cool_gray",
                    )


def build_title_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb("bg_1")
    add_top_rule(slide, prs.slide_width)

    right_bg = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.OVAL,
        Inches(10.2),
        Inches(-0.3),
        Inches(3.8),
        Inches(3.8),
    )
    right_bg.fill.solid()
    right_bg.fill.fore_color.rgb = rgb("bg_2")
    remove_line(right_bg)

    right_accent = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.OVAL,
        Inches(11.05),
        Inches(0.42),
        Inches(2.0),
        Inches(2.0),
    )
    right_accent.fill.solid()
    right_accent.fill.fore_color.rgb = rgb("accent")
    remove_line(right_accent)

    bottom_oval = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.OVAL,
        Inches(9.1),
        Inches(5.1),
        Inches(2.8),
        Inches(1.55),
    )
    bottom_oval.fill.solid()
    bottom_oval.fill.fore_color.rgb = rgb("cool_gray")
    remove_line(bottom_oval)

    add_label(
        slide,
        Inches(0.92),
        Inches(0.75),
        Inches(3.0),
        Inches(0.35),
        "HACK-A-BOT 2026",
        size=10,
        color="accent_soft",
        bold=True,
    )

    title = slide.shapes.title
    title.text = "Hack-A-Bot 2026"
    style_paragraph(title.text_frame.paragraphs[0], font_name="Space Grotesk", size=28, color="text", bold=True)

    subtitle = slide.placeholders[1]
    subtitle.text = "Google Slides theme deck built from the documented site palette and visual board."
    style_paragraph(subtitle.text_frame.paragraphs[0], font_name="Arial", size=14, color="text_dim")

    add_chip(
        slide,
        Inches(0.92),
        Inches(4.0),
        Inches(1.8),
        Inches(0.38),
        "Dark canvas",
        fill="bg_2",
        text_color="text",
        line="line",
    )
    add_chip(
        slide,
        Inches(2.88),
        Inches(4.0),
        Inches(2.3),
        Inches(0.38),
        "Cool-neutral frame",
        fill="bg_2",
        text_color="text",
        line="line",
    )
    add_chip(
        slide,
        Inches(5.34),
        Inches(4.0),
        Inches(2.3),
        Inches(0.38),
        "Red interaction accent",
        fill="accent",
        text_color="neutral",
    )

    add_label(
        slide,
        Inches(0.92),
        Inches(4.72),
        Inches(6.8),
        Inches(0.8),
        "Import this PPTX into Google Slides, then keep it as a starter deck or use it as the theme source for a new presentation.",
        size=14,
        color="text_soft",
        font_name="Arial",
    )


def build_principles_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    add_top_rule(slide, prs.slide_width)
    slide.shapes.title.text = "Theme Principles"
    style_paragraph(slide.shapes.title.text_frame.paragraphs[0], font_name="Space Grotesk", size=22, color="text", bold=True)

    add_panel(slide, Inches(0.72), Inches(1.62), Inches(11.15), Inches(4.95), fill="panel", transparency=0.82, line="cool_gray")

    body = slide.placeholders[1].text_frame
    body.clear()
    body.word_wrap = True
    items = [
        "Dark technical base with cool-gray structure.",
        "High-contrast type for headings, agendas, and scorecards.",
        "Red reserved for action states, emphasis, and key data points.",
        "Borders and glows stay subtle so content remains primary.",
    ]
    for index, item in enumerate(items):
        paragraph = body.paragraphs[0] if index == 0 else body.add_paragraph()
        paragraph.text = item
        paragraph.level = 0
        paragraph.space_after = Pt(12)
        style_paragraph(paragraph, font_name="Arial", size=20, color="text_dim")

    add_chip(
        slide,
        Inches(9.45),
        Inches(0.56),
        Inches(2.05),
        Inches(0.36),
        "Theme Source Deck",
        fill="accent_strong",
        text_color="neutral",
    )


def build_section_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    add_top_rule(slide, prs.slide_width)
    slide.placeholders[0].text = "Event Structure"
    style_paragraph(slide.placeholders[0].text_frame.paragraphs[0], font_name="Space Grotesk", size=26, color="text", bold=True)
    slide.placeholders[1].text = "Use this layout for agenda breaks, tracks, judging phases, and sponsor sections."
    style_paragraph(slide.placeholders[1].text_frame.paragraphs[0], font_name="Arial", size=14, color="text_dim")

    bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0.85),
        Inches(2.2),
        Inches(2.35),
        Inches(0.16),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = rgb("accent")
    remove_line(bar)

    add_label(
        slide,
        Inches(0.85),
        Inches(2.55),
        Inches(2.5),
        Inches(0.55),
        "03",
        size=34,
        color="accent_soft",
        bold=True,
    )


def build_tracks_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[3])
    add_top_rule(slide, prs.slide_width)
    slide.shapes.title.text = "Tracks and Deliverables"
    style_paragraph(slide.shapes.title.text_frame.paragraphs[0], font_name="Space Grotesk", size=22, color="text", bold=True)

    add_panel(slide, Inches(0.78), Inches(1.68), Inches(5.35), Inches(4.75), fill="panel", transparency=0.8, line="cool_gray")
    add_panel(slide, Inches(6.9), Inches(1.68), Inches(5.35), Inches(4.75), fill="panel", transparency=0.8, line="cool_gray")

    left = slide.placeholders[1].text_frame
    left.clear()
    for index, item in enumerate(["AI assistants", "Automation tools", "Student workflows", "Team operations"]):
        paragraph = left.paragraphs[0] if index == 0 else left.add_paragraph()
        paragraph.text = item
        paragraph.level = 0
        paragraph.space_after = Pt(10)
        style_paragraph(paragraph, font_name="Arial", size=18, color="text_dim")

    right = slide.placeholders[2].text_frame
    right.clear()
    for index, item in enumerate(["Demo-ready prototype", "Pitch narrative", "Judging summary", "Deployment next step"]):
        paragraph = right.paragraphs[0] if index == 0 else right.add_paragraph()
        paragraph.text = item
        paragraph.level = 0
        paragraph.space_after = Pt(10)
        style_paragraph(paragraph, font_name="Arial", size=18, color="text_dim")

    add_chip(
        slide,
        Inches(0.95),
        Inches(1.28),
        Inches(1.6),
        Inches(0.34),
        "Tracks",
        fill="bg_2",
        text_color="text",
        line="line",
    )
    add_chip(
        slide,
        Inches(7.08),
        Inches(1.28),
        Inches(1.95),
        Inches(0.34),
        "Deliverables",
        fill="accent",
        text_color="neutral",
    )


def build_metrics_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    add_top_rule(slide, prs.slide_width)
    slide.shapes.title.text = "Presentation Blocks"
    style_paragraph(slide.shapes.title.text_frame.paragraphs[0], font_name="Space Grotesk", size=22, color="text", bold=True)

    card_specs = [
        (Inches(0.78), "48", "Teams"),
        (Inches(4.45), "12", "Hours"),
        (Inches(8.12), "6", "Judges"),
    ]
    for left, value, label in card_specs:
        add_panel(slide, left, Inches(2.0), Inches(3.0), Inches(2.8), fill="panel", transparency=0.78, line="cool_gray")
        add_label(slide, left + Inches(0.28), Inches(2.42), Inches(2.3), Inches(0.7), value, size=30, color="accent", bold=True)
        add_label(slide, left + Inches(0.3), Inches(3.2), Inches(2.1), Inches(0.4), label, size=13, color="text_dim", font_name="Arial")
        add_label(
            slide,
            left + Inches(0.3),
            Inches(3.72),
            Inches(2.3),
            Inches(0.5),
            "Use for metrics, counts, or deadlines.",
            size=11,
            color="text_soft",
            font_name="Arial",
        )


def build_palette_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb("bg_1")
    add_top_rule(slide, prs.slide_width)

    add_label(slide, Inches(0.75), Inches(0.55), Inches(4.2), Inches(0.5), "Palette Reference", size=22, color="text", bold=True)
    add_label(
        slide,
        Inches(0.75),
        Inches(1.0),
        Inches(4.8),
        Inches(0.5),
        "Core colors taken from the documented site theme.",
        size=12,
        color="text_dim",
        font_name="Arial",
    )

    swatches = [
        ("bg_1", "Canvas"),
        ("bg_2", "Section"),
        ("text", "Text"),
        ("text_dim", "Secondary"),
        ("accent", "Accent"),
        ("accent_strong", "Hover"),
        ("cool_gray", "Cool gray"),
        ("neutral", "Neutral"),
    ]
    left = Inches(0.8)
    top = Inches(1.85)
    for index, (color_key, label) in enumerate(swatches):
        col = index % 2
        row = index // 2
        x = left + Inches(2.35 * col)
        y = top + Inches(1.1 * row)
        swatch = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            x,
            y,
            Inches(0.9),
            Inches(0.62),
        )
        swatch.fill.solid()
        swatch.fill.fore_color.rgb = rgb(color_key)
        swatch.line.color.rgb = rgb("cool_gray")
        add_label(slide, x + Inches(1.08), y + Inches(0.05), Inches(1.1), Inches(0.25), label, size=11, color="text", bold=True)
        add_label(
            slide,
            x + Inches(1.08),
            y + Inches(0.3),
            Inches(1.2),
            Inches(0.2),
            f"#{PALETTE[color_key]}",
            size=10,
            color="text_dim",
            font_name="Arial",
        )

    if REFERENCE_IMAGE.exists():
        slide.shapes.add_picture(str(REFERENCE_IMAGE), Inches(5.65), Inches(1.1), width=Inches(6.85))


def patch_theme_xml(pptx_path: Path) -> None:
    with NamedTemporaryFile(suffix=".pptx", delete=False, dir=DOCS_DIR) as temp_file:
        temp_path = Path(temp_file.name)

    color_map = {
        "dk1": PALETTE["bg_1"],
        "lt1": PALETTE["text"],
        "dk2": PALETTE["bg_2"],
        "lt2": PALETTE["neutral"],
        "accent1": PALETTE["accent"],
        "accent2": PALETTE["accent_strong"],
        "accent3": PALETTE["line"],
        "accent4": PALETTE["cool_gray"],
        "accent5": PALETTE["text_dim"],
        "accent6": PALETTE["accent_soft"],
        "hlink": PALETTE["accent"],
        "folHlink": PALETTE["accent_strong"],
    }

    with zipfile.ZipFile(pptx_path) as source, zipfile.ZipFile(temp_path, "w") as target:
        for item in source.infolist():
            data = source.read(item.filename)
            if item.filename == "ppt/theme/theme1.xml":
                root = etree.fromstring(data)
                color_scheme = root.find(".//a:clrScheme", namespaces=NS)
                color_scheme.set("name", "HackABot 2026")
                for key, value in color_map.items():
                    node = color_scheme.find(f"a:{key}", namespaces=NS)
                    for child in list(node):
                        node.remove(child)
                    etree.SubElement(node, f"{{{NS['a']}}}srgbClr", val=value)

                major_latin = root.find(".//a:majorFont/a:latin", namespaces=NS)
                minor_latin = root.find(".//a:minorFont/a:latin", namespaces=NS)
                major_latin.set("typeface", "Space Grotesk")
                minor_latin.set("typeface", "Arial")
                data = etree.tostring(root, xml_declaration=True, encoding="UTF-8", standalone="yes")
            target.writestr(item, data)

    temp_path.replace(pptx_path)


def build_presentation() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    style_layouts(prs)
    build_title_slide(prs)
    build_principles_slide(prs)
    build_section_slide(prs)
    build_tracks_slide(prs)
    build_metrics_slide(prs)
    build_palette_slide(prs)

    prs.save(OUTPUT_PATH)
    patch_theme_xml(OUTPUT_PATH)
    return OUTPUT_PATH


if __name__ == "__main__":
    output = build_presentation()
    print(output)
